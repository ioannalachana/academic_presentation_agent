"""
Academic Presentation Agent — REAL AGENT VERSION
-------------------------------------------------
This is a proper agent, not just a pipeline. It:
  1. PLANS    — analyses the paper and decides a strategy
  2. GENERATES — creates slides and poster
  3. REVIEWS  — scores its own output (0-10)
  4. RETRIES  — rewrites if score < 8, up to 3 attempts
  5. BUILDS   — produces the final .pptx

POST /generate  → runs the full agent loop, returns JSON
GET  /download  → returns the .pptx file
GET  /log       → returns live agent thought log
"""

import os, json, time, uuid, tempfile
from collections import defaultdict

import fitz
from google import genai

from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── App ────────────────────────────────────────────────────────────────────────
app = FastAPI(title="Academic Presentation Agent")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])
app.mount("/static", StaticFiles(directory="static"), name="static")

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")
client = genai.Client(api_key=GEMINI_API_KEY)
MODEL  = "gemini-2.5-flash"

PPTX_STORE: dict = {}
LOG_STORE:  dict = defaultdict(list)

# ── Gemini call with retry ─────────────────────────────────────────────────────
def call_gemini(prompt: str, retries: int = 3) -> str:
    for attempt in range(retries):
        try:
            return client.models.generate_content(model=MODEL, contents=prompt).text.strip()
        except Exception as e:
            if "503" in str(e) or "UNAVAILABLE" in str(e):
                time.sleep(10 * (attempt + 1))
            else:
                raise
    raise HTTPException(503, "Gemini is busy. Please try again in a minute.")

def parse_json(raw: str) -> dict:
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    return json.loads(raw.strip())

def log(sid: str, stage: str, message: str, data: dict = None):
    entry = {"stage": stage, "message": message, "data": data or {}}
    LOG_STORE[sid].append(entry)
    print(f"[{sid[:8]}] [{stage}] {message}")

# ══════════════════════════════════════════════════════
#  AGENT TOOLS
# ══════════════════════════════════════════════════════

def tool_extract_text(pdf_bytes: bytes) -> str:
    doc  = fitz.open(stream=pdf_bytes, filetype="pdf")
    text = "\n\n".join(page.get_text() for page in doc)
    doc.close()
    return text[:100_000]

def tool_plan(sid: str, paper_text: str) -> dict:
    log(sid, "PLAN", "Analysing paper to form a strategy...")
    prompt = f"""You are an expert academic presentation strategist.
Analyse this research paper and return a strategy plan as JSON only, no markdown:

{{
  "paper_type": "empirical|theoretical|review|case_study|methodology",
  "domain": "the academic field",
  "length": "short|medium|long",
  "key_contribution": "one sentence — the main novel contribution",
  "recommended_slides": <number 10-15>,
  "emphasis": ["sections to emphasise e.g. Results, Methods"],
  "audience": "specialist|mixed|general",
  "tone": "technical|accessible|balanced",
  "strategy_notes": "2-3 sentences on the presentation approach"
}}

PAPER (first 8000 chars):
{paper_text[:8000]}"""
    plan = parse_json(call_gemini(prompt))
    log(sid, "PLAN", f"Strategy ready — {plan.get('domain')} / {plan.get('paper_type')} / audience: {plan.get('audience')}", plan)
    return plan

def tool_generate_slides(sid: str, paper_text: str, plan: dict, attempt: int = 1) -> dict:
    log(sid, "GENERATE", f"Generating slide outline (attempt {attempt})...")
    improvements = plan.get("_slide_improvements", [])
    fix_note = ""
    if improvements:
        fix_note = f"\n\nCRITICAL — fix these issues from the previous attempt:\n" + "\n".join(f"- {i}" for i in improvements)
    prompt = f"""You are an expert academic presentation designer.

STRATEGY: paper_type={plan.get('paper_type')}, domain={plan.get('domain')},
key_contribution={plan.get('key_contribution')}, slides={plan.get('recommended_slides')},
emphasise={plan.get('emphasis')}, audience={plan.get('audience')}, tone={plan.get('tone')}
{fix_note}

Create a slide deck as JSON only, no markdown:
{{
  "title": "compelling title",
  "subtitle": "conference or subtitle",
  "total_slides": <number>,
  "slides": [{{
    "slide_number": 1,
    "title": "slide title",
    "type": "title|intro|background|methods|results|discussion|conclusion|references",
    "key_points": ["concise bullet max 15 words"],
    "speaker_notes": "2-4 natural presenter sentences"
  }}]
}}

PAPER:
{paper_text}"""
    data = parse_json(call_gemini(prompt))
    log(sid, "GENERATE", f"Generated {data.get('total_slides')} slides")
    return data

def tool_review_slides(sid: str, slides_data: dict, plan: dict) -> dict:
    log(sid, "REVIEW", "Reviewing slide quality...")
    prompt = f"""You are a strict academic presentation reviewer. Score this slide deck.
Return JSON only, no markdown:
{{
  "score": <1-10>,
  "passes": <true if score >= 8>,
  "issues": ["specific issue"],
  "improvements": ["specific fix"],
  "verdict": "one sentence"
}}

Criteria: key contribution clear (3pts), bullets concise/specific (2pts),
logical narrative arc (2pts), useful speaker notes (2pts), slide count appropriate (1pt).

PLAN: {json.dumps(plan)}
SLIDES: {json.dumps(slides_data)}"""
    review = parse_json(call_gemini(prompt))
    log(sid, "REVIEW", f"Slides score: {review.get('score')}/10 — {'PASS' if review.get('passes') else 'FAIL — rewriting'}", review)
    return review

def tool_generate_poster(sid: str, paper_text: str, plan: dict, attempt: int = 1) -> dict:
    log(sid, "GENERATE", f"Generating poster layout (attempt {attempt})...")
    improvements = plan.get("_poster_improvements", [])
    fix_note = ""
    if improvements:
        fix_note = f"\n\nCRITICAL — fix these issues from the previous attempt:\n" + "\n".join(f"- {i}" for i in improvements)
    prompt = f"""You are an expert academic conference poster designer.

STRATEGY: domain={plan.get('domain')}, key_contribution={plan.get('key_contribution')},
audience={plan.get('audience')}, emphasise={plan.get('emphasis')}
{fix_note}

Create a poster layout as JSON only, no markdown:
{{
  "title": "poster title",
  "authors": "names and affiliations",
  "sections": [{{
    "name": "section name",
    "column": "left|center|right",
    "content": "2-5 clear sentences",
    "bullet_points": ["bullet"],
    "figures": ["figure description or null"]
  }}],
  "key_takeaway": "single most important finding — bold and memorable",
  "contact": "email@university.edu"
}}

Guidelines: 3-column A0, 2-3 sections per column, scannable in 30 seconds.

PAPER:
{paper_text}"""
    data = parse_json(call_gemini(prompt))
    log(sid, "GENERATE", f"Generated {len(data.get('sections', []))} poster sections")
    return data

def tool_review_poster(sid: str, poster_data: dict, plan: dict) -> dict:
    log(sid, "REVIEW", "Reviewing poster quality...")
    prompt = f"""You are a strict academic poster reviewer. Score this poster.
Return JSON only, no markdown:
{{
  "score": <1-10>,
  "passes": <true if score >= 8>,
  "issues": ["specific issue"],
  "improvements": ["specific fix"],
  "verdict": "one sentence"
}}

Criteria: balanced columns (2pts), memorable key_takeaway (2pts),
scannable content (2pts), reflects main contribution (2pts), figures suggested (2pts).

PLAN: {json.dumps(plan)}
POSTER: {json.dumps(poster_data)}"""
    review = parse_json(call_gemini(prompt))
    log(sid, "REVIEW", f"Poster score: {review.get('score')}/10 — {'PASS' if review.get('passes') else 'FAIL — rewriting'}", review)
    return review

def tool_build_pptx(slides_data: dict) -> str:
    DARK_BG = RGBColor(0x1E, 0x27, 0x61)
    LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
    ACCENT = RGBColor(0x00, 0x8B, 0xCC)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
    MID_TEXT = RGBColor(0x47, 0x5A, 0x6E)
    DARK_TYPES = {"title", "conclusion", "references"}

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for info in slides_data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        stype = info.get("type", "content")
        is_dark = stype in DARK_TYPES
        bg_col = DARK_BG if is_dark else LIGHT_BG
        body_col = RGBColor(0xCA, 0xDC, 0xFC) if is_dark else DARK_TEXT

        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = bg_col; bg.line.fill.background()

        nb = slide.shapes.add_textbox(Inches(8.8), Inches(5.1), Inches(1.0), Inches(0.35))
        nb.text_frame.text = str(info["slide_number"])
        nr = nb.text_frame.paragraphs[0]; nr.alignment = PP_ALIGN.RIGHT
        nr.runs[0].font.size = Pt(10)
        nr.runs[0].font.color.rgb = MID_TEXT if not is_dark else RGBColor(0x8A, 0x9B, 0xAF)

        if stype == "title":
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.8), Inches(1.8))
            tb.text_frame.word_wrap = True; tb.text_frame.text = info["title"]
            p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.bold = True; p.runs[0].font.size = Pt(36); p.runs[0].font.color.rgb = WHITE
            kp = info.get("key_points", [])
            if kp:
                sb = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(8.0), Inches(1.0))
                sb.text_frame.text = kp[0]
                sp = sb.text_frame.paragraphs[0]; sp.alignment = PP_ALIGN.CENTER
                sp.runs[0].font.size = Pt(18); sp.runs[0].font.italic = True
                sp.runs[0].font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC)
        else:
            bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
            bar.fill.solid(); bar.fill.fore_color.rgb = DARK_BG if is_dark else ACCENT; bar.line.fill.background()
            ttb = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.8))
            ttb.text_frame.word_wrap = True; ttb.text_frame.text = info["title"]
            tp = ttb.text_frame.paragraphs[0]; tp.alignment = PP_ALIGN.LEFT
            tp.runs[0].font.bold = True; tp.runs[0].font.size = Pt(26); tp.runs[0].font.color.rgb = WHITE
            kp = info.get("key_points", [])
            if kp:
                cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.9))
                cb.text_frame.word_wrap = True
                for i, point in enumerate(kp):
                    p = cb.text_frame.paragraphs[0] if i == 0 else cb.text_frame.add_paragraph()
                    p.text = f"• {point}"; p.space_after = Pt(8)
                    p.runs[0].font.size = Pt(16); p.runs[0].font.color.rgb = body_col

        notes = info.get("speaker_notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4()}.pptx")
    prs.save(path)
    return path

# ══════════════════════════════════════════════════════
#  AGENT LOOP
# ══════════════════════════════════════════════════════

def run_agent(sid: str, pdf_bytes: bytes) -> dict:
    """
    The real agent loop — plans, generates, reviews, retries if needed.
    """
    MAX_ATTEMPTS = 3
    log(sid, "START", "Agent starting...")

    # Step 1: Extract
    paper_text = tool_extract_text(pdf_bytes)
    log(sid, "EXTRACT", f"Extracted {len(paper_text):,} characters from PDF")
    if len(paper_text) < 200:
        raise HTTPException(400, "Could not extract text from this PDF.")

    # Step 2: Plan (agent thinks before acting)
    plan = tool_plan(sid, paper_text)

    # Step 3: Slides — generate → review → retry loop
    slides_data = slides_review = None
    for attempt in range(1, MAX_ATTEMPTS + 1):
        slides_data   = tool_generate_slides(sid, paper_text, plan, attempt)
        slides_review = tool_review_slides(sid, slides_data, plan)
        if slides_review.get("passes"):
            log(sid, "AGENT", f"Slides accepted after {attempt} attempt(s) ✓")
            break
        if attempt < MAX_ATTEMPTS:
            plan["_slide_improvements"] = slides_review.get("improvements", [])
        else:
            log(sid, "AGENT", "Max attempts reached — using best slides available")

    # Step 4: Poster — generate → review → retry loop
    poster_data = poster_review = None
    for attempt in range(1, MAX_ATTEMPTS + 1):
        poster_data   = tool_generate_poster(sid, paper_text, plan, attempt)
        poster_review = tool_review_poster(sid, poster_data, plan)
        if poster_review.get("passes"):
            log(sid, "AGENT", f"Poster accepted after {attempt} attempt(s) ✓")
            break
        if attempt < MAX_ATTEMPTS:
            plan["_poster_improvements"] = poster_review.get("improvements", [])
        else:
            log(sid, "AGENT", "Max attempts reached — using best poster available")

    # Step 5: Build .pptx
    log(sid, "BUILD", "Building PowerPoint file...")
    PPTX_STORE[sid] = tool_build_pptx(slides_data)
    log(sid, "DONE", "Agent complete ✓")

    return {
        "session_id":    sid,
        "plan":          plan,
        "slides":        slides_data,
        "slides_review": slides_review,
        "poster":        poster_data,
        "poster_review": poster_review,
        "agent_log":     LOG_STORE[sid],
    }

# ── Routes ─────────────────────────────────────────────────────────────────────

@app.get("/", response_class=HTMLResponse)
async def index():
    with open("static/index.html") as f:
        return f.read()

@app.post("/generate")
async def generate(pdf: UploadFile = File(...)):
    if not pdf.filename.lower().endswith(".pdf"):
        raise HTTPException(400, "Please upload a PDF file.")
    pdf_bytes = await pdf.read()
    sid = str(uuid.uuid4())
    return run_agent(sid, pdf_bytes)

@app.get("/download/{session_id}")
async def download(session_id: str):
    path = PPTX_STORE.get(session_id)
    if not path or not os.path.exists(path):
        raise HTTPException(404, "File not found. Please generate again.")
    return FileResponse(path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="presentation.pptx")

@app.get("/log/{session_id}")
async def get_log(session_id: str):
    return {"log": LOG_STORE.get(session_id, [])}
